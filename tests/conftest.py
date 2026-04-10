# ==============================================================================
# Informity AI — Test Configuration
# Shared fixtures for all tests.
# ==============================================================================

import sys
from pathlib import Path
from unittest.mock import patch

import docx
import openpyxl
import pptx
import pymupdf
import pytest

# ==============================================================================
# Fixture Directory
# ==============================================================================

# Ensure repository-root imports (e.g., tools.diagnostics.*) resolve reliably
# in CI regardless of how pytest is invoked.
_REPO_ROOT = Path(__file__).resolve().parents[1]
if str(_REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(_REPO_ROOT))


@pytest.fixture(autouse=True)
def _disable_reranker_in_tests():
    # Prevent the real cross-encoder reranker from loading during tests.
    # Patch only rag_rerank on the real settings object so that:
    # - Real paths (cache_dir, models_dir, etc.) remain intact for docling/embedder
    # - configure_hf_environment() and are_required_models_cached() work correctly
    # - Docling extractor tests pass when models are cached (e.g. after install)
    try:
        from informity.indexer.adaptive_tuning import invalidate_tuning_cache
        invalidate_tuning_cache()
    except ImportError:
        pass
    from informity.config import settings as real_settings
    with patch.object(real_settings, 'rag_rerank', False):
        yield


# ==============================================================================
# Sample File Generators
# ==============================================================================


@pytest.fixture
def sample_txt(tmp_path: Path) -> Path:
    # Create a sample .txt file.
    f = tmp_path / "sample.txt"
    f.write_text(
        "Hello, Informity AI!\nThis is a test text file.\nWith three lines.\n", encoding="utf-8"
    )
    return f


@pytest.fixture
def sample_md(tmp_path: Path) -> Path:
    # Create a sample .md file.
    f = tmp_path / "sample.md"
    f.write_text(
        "# Heading\n\nA paragraph with **bold** text.\n\n- Item 1\n- Item 2\n", encoding="utf-8"
    )
    return f


@pytest.fixture
def sample_pdf(tmp_path: Path) -> Path:
    # Create a sample .pdf file with two pages.
    f = tmp_path / "sample.pdf"
    doc = pymupdf.open()
    page1 = doc.new_page()
    page1.insert_text((72, 72), "Page one of the test PDF.\nInformity AI document.")
    page2 = doc.new_page()
    page2.insert_text((72, 72), "Page two with additional content.")
    doc.save(str(f))
    doc.close()
    return f


@pytest.fixture
def sample_docx(tmp_path: Path) -> Path:
    # Create a sample .docx file with paragraphs and a table.
    f = tmp_path / "sample.docx"
    doc = docx.Document()
    doc.core_properties.title = "Test DOCX"
    doc.core_properties.author = "Informity AI"
    doc.add_heading("Document Title", level=1)
    doc.add_paragraph("First paragraph of the test document.")
    doc.add_paragraph("Second paragraph with more text.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Col A"
    table.cell(0, 1).text = "Col B"
    table.cell(1, 0).text = "Val 1"
    table.cell(1, 1).text = "Val 2"
    doc.save(str(f))
    return f


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    # Create a sample .pptx file with two slides.
    f = tmp_path / "sample.pptx"
    prs = pptx.Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "Slide One Title"
    slide1.placeholders[1].text = "Body text on slide one."
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Slide Two Title"
    slide2.placeholders[1].text = "Body text on slide two."
    prs.save(str(f))
    return f


@pytest.fixture
def sample_xlsx(tmp_path: Path) -> Path:
    # Create a sample .xlsx file with data.
    f = tmp_path / "sample.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["Name", "Age", "City"])
    ws.append(["Alice", "30", "NYC"])
    ws.append(["Bob", "25", "LA"])
    wb.save(str(f))
    return f


@pytest.fixture
def sample_csv(tmp_path: Path) -> Path:
    # Create a sample .csv file.
    f = tmp_path / "sample.csv"
    f.write_text("name,age,city\nAlice,30,NYC\nBob,25,LA\n", encoding="utf-8")
    return f


@pytest.fixture
def sample_html(tmp_path: Path) -> Path:
    # Create a sample .html file.
    f = tmp_path / "sample.html"
    f.write_text(
        "<!DOCTYPE html>\n"
        "<html><head><title>Test Page</title></head>\n"
        "<body>\n"
        "<h1>Welcome</h1>\n"
        "<p>A paragraph of text.</p>\n"
        "<script>var x = 1;</script>\n"
        "</body></html>\n",
        encoding="utf-8",
    )
    return f
